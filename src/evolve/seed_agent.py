"""Evolvable agent harness — fully self-contained, single-file AI coding agent.

This file is the "genome" that SkyDiscover evolves. It contains EVERYTHING the
agent needs: data classes, configuration, system prompt, tool declarations,
tool implementations, the agent loop, and helper functions.

NO external imports from this project — only stdlib + google-genai + document libs.
This ensures the LLM evolving this code can see and modify every definition.

INTERFACE CONTRACT (must be preserved for the evaluator to load this agent):
  - Must define AgentResult dataclass with: response, tool_calls, messages, error
  - Must define a class with:
      - async def run(self, prompt: str, cwd: Path) -> AgentResult
      - def name(self) -> str
  - The class name can be anything, but it must have those two methods.
"""

from __future__ import annotations

# ── IMPORTS (stdlib + google-genai only) ────────────────────────────────────

import asyncio
import fnmatch
import logging
import os
import re
import subprocess
import time
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from google import genai
from google.genai.types import (
    Content,
    FunctionDeclaration,
    GenerateContentConfig,
    Part,
    Tool,
)

logger = logging.getLogger(__name__)


# ── DATA CLASSES ────────────────────────────────────────────────────────────

@dataclass
class AgentResult:
    """Result of running the agent on a single task.

    Fields:
        response:   The agent's final text output + any deliverable file contents.
        tool_calls: Log of tool invocations [{tool: str, input: str}, ...].
        messages:   Full conversation log for traces.
        error:      Error message if the agent failed, else None.
    """
    response: str = ""
    tool_calls: list[dict[str, Any]] = field(default_factory=list)
    messages: list[dict[str, Any]] = field(default_factory=list)
    error: str | None = None


class BaseAgent(ABC):
    """Abstract base for agent backends.

    Subclasses must implement:
        - async def run(self, prompt: str, cwd: Path) -> AgentResult
        - def name(self) -> str
    """

    def __init__(
        self,
        system_prompt: str | None = None,
        max_turns: int = 10,
        model: str | None = None,
    ) -> None:
        self._system_prompt = system_prompt
        self._max_turns = max_turns
        self._model = model

    @abstractmethod
    async def run(self, prompt: str, cwd: Path) -> AgentResult:
        """Run the agent on a task prompt in the given working directory."""
        ...

    @abstractmethod
    def name(self) -> str:
        """Human-readable agent name."""
        ...


# ── CONFIGURATION ───────────────────────────────────────────────────────────

DEFAULT_MODEL = "gemini-2.5-flash"
MAX_ITERATIONS = 30          # max tool-call rounds per task
TEMPERATURE = 0.2            # low for deterministic, reliable output
MAX_OUTPUT_TOKENS = 16384    # generous for code generation
BASH_TIMEOUT = 120           # seconds per bash command
MAX_BASH_OUTPUT = 15000      # chars of stdout/stderr to keep
MAX_FILE_READ = 60000        # chars before truncating file reads
MAX_GREP_RESULTS = 50        # max matches returned by grep
MAX_GLOB_RESULTS = 100       # max files returned by glob

# ── SYSTEM PROMPT ───────────────────────────────────────────────────────────
#
# This is the core instruction set for the LLM. It determines how the agent
# approaches tasks, uses tools, and verifies its work. Changes here have the
# most direct impact on task completion scores.

SYSTEM_PROMPT = """\
You are an expert professional assistant. You complete tasks by writing and \
executing code in the working directory: {cwd}

RULES:
- Work until FULLY done. Never stop early. Never ask questions.
- Fix your own errors. If something fails, read the error and fix it.
- All paths are relative to {cwd}. Never write outside it.
- Existing files are inputs/references. Create NEW output files.

WORKFLOW:
1. EXPLORE — Run list_dir to see what files exist. Read ALL reference files \
with read_file (it natively handles .pdf, .docx, .xlsx, .pptx, .csv, .txt, \
.json). Extract every piece of data you will need.

2. PLAN — Identify EVERY required deliverable. Note exact file names, formats, \
and content requirements. Count them. You will verify this count at the end.

3. CREATE — For each deliverable:
   - Binary formats (.xlsx, .docx, .pptx, .pdf): write a Python script using \
the appropriate library (openpyxl, python-docx, python-pptx, reportlab/fpdf2), \
then run it with bash.
   - Text formats (.csv, .txt, .py, .json, .html, .md): use write_file directly.
   - Include ALL specific data, numbers, names, dates from the task and reference \
files. Do not omit or summarize — include them EXACTLY.

4. VERIFY — After creating all files:
   - Run list_dir to confirm every deliverable exists.
   - Run read_file on each deliverable to check its content matches requirements.
   - If anything is wrong or missing, fix it now.

TOOL GUIDE:
- bash: run scripts, install packages, system commands. Check exit code.
- read_file: read any file. Handles .pdf, .docx, .xlsx, .pptx natively. \
Use offset/limit for large files.
- write_file: create or overwrite text files.
- edit_file: surgical find-and-replace in existing files. old_string must match exactly.
- list_dir: see directory contents with file sizes.
- grep: search file contents by regex.
- glob: find files by name pattern (supports **).

COMMON MISTAKES TO AVOID:
- Reading reference files but forgetting to use their data in deliverables.
- Creating only some deliverables and stopping.
- Using .csv when .xlsx was requested (or similar format mismatches).
- Not verifying that output files actually contain the required content.
- Writing Python scripts that import libraries without installing them first.
"""

# ── TOOL DEFINITIONS ────────────────────────────────────────────────────────

TOOL_DECLARATIONS = Tool(function_declarations=[
    FunctionDeclaration(
        name="bash",
        description=(
            "Execute a shell command. Returns stdout, stderr, and exit code. "
            "Use for: running Python scripts, installing packages (pip install), "
            "file operations, and any system commands. Always check the exit code."
        ),
        parameters={
            "type": "object",
            "properties": {
                "command": {
                    "type": "string",
                    "description": "The shell command to execute.",
                },
            },
            "required": ["command"],
        },
    ),
    FunctionDeclaration(
        name="read_file",
        description=(
            "Read file contents. Handles all common formats natively:\n"
            "- .pdf → extracts text from all pages\n"
            "- .docx → extracts paragraphs and tables\n"
            "- .xlsx/.xlsm → reads all sheets with cell data\n"
            "- .pptx → extracts slide text\n"
            "- text files → returns content with line numbers\n"
            "Use offset and limit for large files."
        ),
        parameters={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "File path relative to working directory.",
                },
                "offset": {
                    "type": "integer",
                    "description": "Start line (1-based). Optional.",
                },
                "limit": {
                    "type": "integer",
                    "description": "Max lines to read. Optional.",
                },
            },
            "required": ["path"],
        },
    ),
    FunctionDeclaration(
        name="write_file",
        description=(
            "Create or overwrite a file with the given content. "
            "Best for text files (.py, .csv, .txt, .json, .html, .md). "
            "For binary formats (.docx, .xlsx), write a Python script and run via bash."
        ),
        parameters={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "File path relative to working directory.",
                },
                "content": {
                    "type": "string",
                    "description": "Complete file content.",
                },
            },
            "required": ["path", "content"],
        },
    ),
    FunctionDeclaration(
        name="edit_file",
        description=(
            "Find-and-replace edit in an existing file. "
            "old_string must match exactly (including whitespace and indentation). "
            "Use for surgical edits; use write_file for full rewrites."
        ),
        parameters={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "File path relative to working directory.",
                },
                "old_string": {
                    "type": "string",
                    "description": "Exact text to find.",
                },
                "new_string": {
                    "type": "string",
                    "description": "Replacement text.",
                },
            },
            "required": ["path", "old_string", "new_string"],
        },
    ),
    FunctionDeclaration(
        name="list_dir",
        description="List files and directories with sizes. Use to verify deliverables exist.",
        parameters={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Directory path. Defaults to '.'.",
                },
            },
        },
    ),
    FunctionDeclaration(
        name="grep",
        description="Search file contents by regex. Returns matching lines with file:line: prefix.",
        parameters={
            "type": "object",
            "properties": {
                "pattern": {
                    "type": "string",
                    "description": "Regex pattern to search for.",
                },
                "path": {
                    "type": "string",
                    "description": "File or directory to search. Defaults to '.'.",
                },
                "include": {
                    "type": "string",
                    "description": "Glob filter for files (e.g. '*.py'). Optional.",
                },
            },
            "required": ["pattern"],
        },
    ),
    FunctionDeclaration(
        name="glob",
        description="Find files by name pattern. Supports ** for recursion. Sorted by modification time.",
        parameters={
            "type": "object",
            "properties": {
                "pattern": {
                    "type": "string",
                    "description": "Glob pattern (e.g. '**/*.py', '*.xlsx').",
                },
            },
            "required": ["pattern"],
        },
    ),
])


# ── TOOL IMPLEMENTATIONS ───────────────────────────────────────────────────

def _check_path(path: str, cwd: Path) -> tuple[Path, str | None]:
    """Resolve path and verify it's within the workspace."""
    resolved = (cwd / path).resolve()
    if not str(resolved).startswith(str(cwd)):
        return resolved, f"Error: path '{path}' is outside the working directory."
    return resolved, None


def _run_bash(command: str, cwd: Path) -> str:
    """Execute a shell command in the workspace."""
    try:
        result = subprocess.run(
            command,
            shell=True,
            cwd=str(cwd),
            capture_output=True,
            text=True,
            timeout=BASH_TIMEOUT,
            env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1"},
        )
        stdout = result.stdout[:MAX_BASH_OUTPUT] if result.stdout else ""
        stderr = result.stderr[:MAX_BASH_OUTPUT] if result.stderr else ""
        parts = [f"Exit code: {result.returncode}"]
        if stdout.strip():
            parts.append(f"STDOUT:\n{stdout}")
        if stderr.strip():
            parts.append(f"STDERR:\n{stderr}")
        return "\n".join(parts)
    except subprocess.TimeoutExpired:
        return f"Command timed out after {BASH_TIMEOUT}s. Try a simpler command."
    except Exception as e:
        return f"Error executing command: {e}"


def _read_pdf(resolved: Path, display_path: str) -> str:
    """Extract text from a PDF using pymupdf."""
    try:
        import pymupdf
        doc = pymupdf.open(str(resolved))
        parts = []
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if text:
                parts.append(f"--- Page {i + 1} ---\n{text}")
        doc.close()
        result = "\n\n".join(parts) if parts else "(PDF has no extractable text)"
        if len(result) > MAX_FILE_READ:
            result = result[:MAX_FILE_READ] + "\n... (truncated)"
        return f"File: {display_path} (PDF, {len(parts)} pages)\n{result}"
    except Exception as e:
        return f"Error reading PDF '{display_path}': {e}"


def _read_office_file(resolved: Path, display_path: str) -> str:
    """Extract text from Office documents (.docx, .xlsx, .pptx)."""
    suffix = resolved.suffix.lower()
    try:
        if suffix == ".docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(str(resolved))
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
            result = "\n".join(parts)
        elif suffix in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(str(resolved), data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    row_str = " | ".join(str(c) if c is not None else "" for c in row)
                    parts.append(row_str)
            result = "\n".join(parts)
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(resolved))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"[Slide {i + 1}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            result = "\n".join(parts)
        else:
            return f"Unsupported file type: {suffix}"
        if len(result) > MAX_FILE_READ:
            result = result[:MAX_FILE_READ] + "\n... (truncated)"
        return f"File: {display_path} ({suffix})\n{result}"
    except Exception as e:
        return f"Error reading '{display_path}': {e}"


def _read_file(path: str, cwd: Path, offset: int = 0, limit: int = 0) -> str:
    """Read a file with native support for PDF, Office docs, and text files."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    if not resolved.exists():
        available = _list_available_files(cwd)
        return f"Error: file '{path}' not found.\nAvailable files:\n{available}"
    if resolved.is_dir():
        return f"Error: '{path}' is a directory. Use list_dir instead."

    if resolved.suffix.lower() == ".pdf":
        return _read_pdf(resolved, path)
    if resolved.suffix.lower() in (".docx", ".xlsx", ".xlsm", ".pptx"):
        return _read_office_file(resolved, path)

    # Text files
    try:
        content = resolved.read_text(encoding="utf-8", errors="strict")
    except (UnicodeDecodeError, ValueError):
        size = resolved.stat().st_size
        return f"Binary file ({size} bytes). Use bash with Python to inspect."

    lines = content.splitlines(keepends=True)
    total = len(lines)
    start = max(0, offset - 1) if offset > 0 else 0
    end = min(start + limit, total) if limit > 0 else total
    selected = lines[start:end]

    numbered = [f"{i:>6}\t{line.rstrip()}" for i, line in enumerate(selected, start=start + 1)]
    result = "\n".join(numbered)
    if len(result) > MAX_FILE_READ:
        result = result[:MAX_FILE_READ] + f"\n... (truncated, {total} total lines)"

    header = f"File: {path} ({total} lines)"
    if start > 0 or end < total:
        header += f" [showing lines {start + 1}-{end}]"
    return f"{header}\n{result}"


def _write_file(path: str, content: str, cwd: Path) -> str:
    """Write content to a file in the workspace."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    try:
        resolved.parent.mkdir(parents=True, exist_ok=True)
        resolved.write_text(content)
        size = resolved.stat().st_size
        return f"Successfully wrote {size} bytes to {path}"
    except Exception as e:
        return f"Error writing file: {e}"


def _edit_file(path: str, old_string: str, new_string: str, cwd: Path) -> str:
    """Find-and-replace edit in a file."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    if not resolved.exists():
        return f"Error: file '{path}' not found."
    try:
        content = resolved.read_text(encoding="utf-8")
    except (UnicodeDecodeError, ValueError):
        return f"Error: '{path}' is binary and cannot be edited."

    count = content.count(old_string)
    if count == 0:
        preview = content[:2000]
        return (
            f"Error: old_string not found in '{path}'. "
            f"Must match exactly (including whitespace).\n"
            f"File preview:\n{preview}"
        )
    if count > 1:
        return f"Error: old_string found {count} times. Add more context to make it unique."

    new_content = content.replace(old_string, new_string, 1)
    resolved.write_text(new_content)

    new_lines = new_content.splitlines()
    edit_start = content.index(old_string)
    line_num = content[:edit_start].count("\n")
    ctx_start = max(0, line_num - 2)
    ctx_end = min(len(new_lines), line_num + new_string.count("\n") + 3)
    context = "\n".join(f"{i + 1:>6}\t{new_lines[i]}" for i in range(ctx_start, ctx_end))
    return f"Successfully edited {path}.\nUpdated region:\n{context}"


def _list_dir(path: str, cwd: Path) -> str:
    """List directory contents with file sizes."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    if not resolved.exists():
        return f"Error: directory '{path}' does not exist."
    if not resolved.is_dir():
        return f"Error: '{path}' is not a directory."
    entries = []
    for item in sorted(resolved.iterdir()):
        if item.name.startswith("."):
            continue
        if item.is_file():
            size = item.stat().st_size
            if size < 1024:
                s = f"{size}B"
            elif size < 1024 * 1024:
                s = f"{size / 1024:.1f}KB"
            else:
                s = f"{size / (1024 * 1024):.1f}MB"
            entries.append(f"  {item.name}  ({s})")
        elif item.is_dir():
            entries.append(f"  {item.name}/")
    return "\n".join(entries) if entries else "(empty directory)"


def _grep(pattern: str, path: str, cwd: Path, include: str = "") -> str:
    """Search file contents using regex."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    if not resolved.exists():
        return f"Error: path '{path}' does not exist."
    try:
        regex = re.compile(pattern)
    except re.error as e:
        return f"Error: invalid regex: {e}"

    matches: list[str] = []

    def _search(fpath: Path) -> None:
        try:
            text = fpath.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return
        rel = str(fpath.relative_to(cwd))
        for i, line in enumerate(text.splitlines(), 1):
            if regex.search(line):
                matches.append(f"{rel}:{i}: {line.rstrip()[:200]}")
                if len(matches) >= MAX_GREP_RESULTS:
                    return

    if resolved.is_file():
        _search(resolved)
    else:
        for fpath in sorted(resolved.rglob("*")):
            if len(matches) >= MAX_GREP_RESULTS:
                break
            if not fpath.is_file() or fpath.name.startswith("."):
                continue
            if include and not fnmatch.fnmatch(fpath.name, include):
                continue
            _search(fpath)

    if not matches:
        return f"No matches for '{pattern}' in {path}."
    header = f"Found {len(matches)} match(es)"
    if len(matches) >= MAX_GREP_RESULTS:
        header += f" (limited to {MAX_GREP_RESULTS})"
    return f"{header}:\n" + "\n".join(matches)


def _glob_files(pattern: str, cwd: Path) -> str:
    """Find files matching a glob pattern."""
    try:
        matches = sorted(cwd.glob(pattern), key=lambda p: p.stat().st_mtime, reverse=True)
    except (OSError, ValueError):
        matches = []
    matches = [m for m in matches if not any(p.startswith(".") for p in m.relative_to(cwd).parts)]
    if not matches:
        return f"No files matching '{pattern}'."
    results = []
    for m in matches[:MAX_GLOB_RESULTS]:
        rel = str(m.relative_to(cwd))
        if m.is_file():
            results.append(f"  {rel}  ({m.stat().st_size} bytes)")
        else:
            results.append(f"  {rel}/")
    header = f"Found {len(matches)} match(es)"
    if len(matches) > MAX_GLOB_RESULTS:
        header += f" (showing first {MAX_GLOB_RESULTS})"
    return f"{header}:\n" + "\n".join(results)


def _list_available_files(cwd: Path) -> str:
    """List workspace root files for error messages."""
    entries = []
    for item in sorted(cwd.iterdir()):
        if item.name.startswith("."):
            continue
        entries.append(f"  {item.name}{'/' if item.is_dir() else ''}")
    return "\n".join(entries) if entries else "(empty directory)"


def _execute_tool(name: str, args: dict[str, Any], cwd: Path) -> str:
    """Dispatch a tool call to its implementation."""
    try:
        if name == "bash":
            return _run_bash(args.get("command", ""), cwd)
        elif name == "read_file":
            return _read_file(
                args.get("path", ""), cwd,
                offset=int(args.get("offset", 0)),
                limit=int(args.get("limit", 0)),
            )
        elif name == "write_file":
            return _write_file(args.get("path", ""), args.get("content", ""), cwd)
        elif name == "edit_file":
            return _edit_file(
                args.get("path", ""), args.get("old_string", ""),
                args.get("new_string", ""), cwd,
            )
        elif name == "list_dir":
            return _list_dir(args.get("path", "."), cwd)
        elif name == "grep":
            return _grep(
                args.get("pattern", ""), args.get("path", "."),
                cwd, include=args.get("include", ""),
            )
        elif name == "glob":
            return _glob_files(args.get("pattern", ""), cwd)
        else:
            return f"Unknown tool: {name}"
    except Exception as e:
        return f"Tool error ({name}): {e}"


# ── AGENT LOOP ──────────────────────────────────────────────────────────────

class CustomAgent(BaseAgent):
    """Self-contained coding agent using Gemini API with tool use."""

    def name(self) -> str:
        return "custom"

    async def run(self, prompt: str, cwd: Path) -> AgentResult:
        """Run the agent. Never raises — errors become partial results."""
        cwd = cwd.resolve()
        cwd.mkdir(parents=True, exist_ok=True)

        tool_calls_log: list[dict[str, Any]] = []
        messages_log: list[dict[str, Any]] = []
        response_parts: list[str] = []

        try:
            return await asyncio.to_thread(
                self._run_inner, prompt, cwd,
                tool_calls_log, messages_log, response_parts,
            )
        except Exception as e:
            logger.error("[custom] Unhandled error: %s: %s", type(e).__name__, e)
            response_parts.append(f"(Agent error: {type(e).__name__}: {e})")
            return AgentResult(
                response="\n".join(response_parts),
                tool_calls=tool_calls_log,
                messages=messages_log,
            )

    def _run_inner(
        self,
        prompt: str,
        cwd: Path,
        tool_calls_log: list[dict[str, Any]],
        messages_log: list[dict[str, Any]],
        response_parts: list[str],
    ) -> AgentResult:
        """Synchronous ReAct loop — runs in a thread for async compatibility."""
        api_key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
        if not api_key:
            return AgentResult(error="GEMINI_API_KEY or GOOGLE_API_KEY required.")
        client = genai.Client(api_key=api_key)

        model = self._model or DEFAULT_MODEL
        max_iters = self._max_turns if self._max_turns != 10 else MAX_ITERATIONS
        system = SYSTEM_PROMPT.format(cwd=cwd)

        config = GenerateContentConfig(
            system_instruction=system,
            temperature=TEMPERATURE,
            max_output_tokens=MAX_OUTPUT_TOKENS,
            tools=[TOOL_DECLARATIONS],
        )

        contents: list[Content | str] = [prompt]
        last_calls: list[str] = []
        nudge_count = 0

        for iteration in range(max_iters):
            logger.debug("[custom] Iteration %d/%d", iteration + 1, max_iters)

            response = _call_with_retry(client, model, contents, config)
            if response is None:
                response_parts.append("Unable to complete task due to API errors.")
                break

            function_calls = _extract_function_calls(response)

            if not function_calls:
                text = _get_response_text(response)
                response_parts.append(text)
                messages_log.append({"role": "assistant", "type": "text", "content": text[:2000]})
                break

            for part in _get_content_parts(response):
                if part.text:
                    response_parts.append(part.text)
                    messages_log.append({"role": "assistant", "type": "text", "content": part.text[:2000]})

            tool_response_parts: list[Part] = []
            for fc in function_calls:
                try:
                    args = dict(fc.args) if fc.args else {}
                except Exception:
                    args = {}
                result = _execute_tool(fc.name, args, cwd)

                tool_calls_log.append({"tool": fc.name, "input": str(args)[:500]})
                messages_log.append({"role": "assistant", "type": "tool_use", "tool": fc.name, "input": str(args)[:1000]})
                messages_log.append({"role": "tool", "type": "tool_result", "content": result[:2000]})

                tool_response_parts.append(
                    Part.from_function_response(name=fc.name, response={"result": result})
                )

                call_sig = f"{fc.name}:{hash(str(args))}"
                last_calls.append(call_sig)

            try:
                model_content = response.candidates[0].content if response.candidates else None
            except (IndexError, AttributeError):
                model_content = None
            if model_content:
                contents.append(model_content)
            contents.append(Content(parts=tool_response_parts))

            # Doom-loop detection: last 3 calls repeat prior 3
            if len(last_calls) >= 6 and last_calls[-3:] == last_calls[-6:-3]:
                nudge_count += 1
                if nudge_count >= 3:
                    response_parts.append("(Agent terminated: stuck in repeated actions)")
                    break
                contents.append(
                    "You are repeating the same actions. Try a different approach or finish."
                )

        return AgentResult(
            response="\n".join(response_parts),
            tool_calls=tool_calls_log,
            messages=messages_log,
        )


# ── HELPERS ─────────────────────────────────────────────────────────────────

def _call_with_retry(
    client: genai.Client,
    model: str,
    contents: list,
    config: GenerateContentConfig,
    max_retries: int = 3,
) -> Any | None:
    """Call Gemini API with exponential backoff."""
    for attempt in range(max_retries + 1):
        try:
            return client.models.generate_content(
                model=model, contents=contents, config=config,
            )
        except Exception as e:
            if attempt == max_retries:
                logger.error("[custom] API failed after %d retries: %s", max_retries, e)
                return None
            wait = 2 ** (attempt + 1)
            logger.warning("[custom] API error (attempt %d): %s — retry in %ds", attempt + 1, e, wait)
            time.sleep(wait)
    return None


def _extract_function_calls(response: Any) -> list:
    """Extract FunctionCall objects from a Gemini response."""
    if not response.candidates:
        return []
    content = response.candidates[0].content
    if not content or not content.parts:
        return []
    return [part.function_call for part in content.parts if part.function_call]


def _get_response_text(response: Any) -> str:
    """Safely extract text from a Gemini response."""
    try:
        return response.text or ""
    except (AttributeError, ValueError):
        if response.candidates:
            content = response.candidates[0].content
            if content and content.parts:
                return "\n".join(p.text for p in content.parts if p.text)
        return ""


def _get_content_parts(response: Any) -> list:
    """Safely get content parts from a Gemini response."""
    if not response.candidates:
        return []
    content = response.candidates[0].content
    return content.parts if content and content.parts else []
