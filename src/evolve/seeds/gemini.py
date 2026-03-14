from __future__ import annotations

import asyncio
import fnmatch
import logging
import os
import re
import subprocess
import time
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from google import genai
from google.genai.types import (
    Content, FunctionDeclaration, GenerateContentConfig, Part, Tool,
)

logger = logging.getLogger(__name__)

# ── Data classes ────────────────────────────────────────────────────────────

@dataclass
class AgentResult:
    """Result of running the agent on a single task."""
    response: str = ""
    tool_calls: list[dict[str, Any]] = field(default_factory=list)
    messages: list[dict[str, Any]] = field(default_factory=list)
    error: str | None = None

class BaseAgent(ABC):
    def __init__(self, system_prompt: str | None = None, max_turns: int = 10, model: str | None = None):
        self._system_prompt = system_prompt
        self._max_turns = max_turns
        self._model = model
    @abstractmethod
    async def run(self, prompt: str, cwd: Path) -> AgentResult: ...
    @abstractmethod
    def name(self) -> str: ...

# ── Configuration ───────────────────────────────────────────────────────────

DEFAULT_MODEL = "gemini-3-flash"
MAX_ITERATIONS = 30
TEMPERATURE = 0.2
MAX_OUTPUT_TOKENS = 16384
BASH_TIMEOUT = 120
MAX_BASH_OUTPUT = 15000
MAX_FILE_READ = 60000
MAX_GREP_RESULTS = 50
MAX_GLOB_RESULTS = 100

# ── System prompt ───────────────────────────────────────────────────────────

SYSTEM_PROMPT = """\
You are an assistant that completes tasks using the provided tools.
Your working directory is: {cwd}

All file paths are relative to the working directory.
Save all output files to the working directory. Never write files outside it.
Keep going until the task is fully complete. Verify your work before finishing.
If a command fails, read the error and fix the issue yourself.
"""

# ── 8 Standard tools — declarations (Gemini format) ────────────────────────

TOOL_DECLARATIONS = Tool(function_declarations=[
    FunctionDeclaration(
        name="bash",
        description="Execute a shell command. Returns stdout, stderr, exit code. Use for running scripts, pip install, system commands.",
        parameters={"type": "object", "properties": {"command": {"type": "string", "description": "Shell command to execute."}}, "required": ["command"]},
    ),
    FunctionDeclaration(
        name="read_file",
        description="Read file contents. Natively handles .pdf, .docx, .xlsx, .pptx, and text files with line numbers. Use offset/limit for large files.",
        parameters={"type": "object", "properties": {"path": {"type": "string", "description": "File path (relative)."}, "offset": {"type": "integer", "description": "Start line (1-based). Optional."}, "limit": {"type": "integer", "description": "Max lines. Optional."}}, "required": ["path"]},
    ),
    FunctionDeclaration(
        name="write_file",
        description="Create or overwrite a file. Best for text files. For .docx/.xlsx, write a Python script and run via bash.",
        parameters={"type": "object", "properties": {"path": {"type": "string", "description": "File path (relative)."}, "content": {"type": "string", "description": "Complete file content."}}, "required": ["path", "content"]},
    ),
    FunctionDeclaration(
        name="edit_file",
        description="Find-and-replace edit. old_string must match exactly including whitespace.",
        parameters={"type": "object", "properties": {"path": {"type": "string", "description": "File path."}, "old_string": {"type": "string", "description": "Exact text to find."}, "new_string": {"type": "string", "description": "Replacement."}}, "required": ["path", "old_string", "new_string"]},
    ),
    FunctionDeclaration(
        name="list_dir",
        description="List directory contents with file sizes.",
        parameters={"type": "object", "properties": {"path": {"type": "string", "description": "Directory path. Defaults to '.'."}}},
    ),
    FunctionDeclaration(
        name="grep",
        description="Search file contents by regex. Returns matching lines with file:line: prefix.",
        parameters={"type": "object", "properties": {"pattern": {"type": "string", "description": "Regex pattern."}, "path": {"type": "string", "description": "Search path. Defaults to '.'."}, "include": {"type": "string", "description": "Glob filter (e.g. '*.py'). Optional."}}, "required": ["pattern"]},
    ),
    FunctionDeclaration(
        name="glob",
        description="Find files by name pattern. Supports ** for recursion.",
        parameters={"type": "object", "properties": {"pattern": {"type": "string", "description": "Glob pattern (e.g. '**/*.py')."}}, "required": ["pattern"]},
    ),
    FunctionDeclaration(
        name="done",
        description="Signal task completion. Call after verifying all deliverables exist and are correct.",
        parameters={"type": "object", "properties": {"summary": {"type": "string", "description": "Brief summary of deliverables created."}}, "required": ["summary"]},
    ),
])

# ── 8 Standard tools — implementations (provider-agnostic) ─────────────────

def _check_path(path: str, cwd: Path) -> tuple[Path, str | None]:
    resolved = (cwd / path).resolve()
    if not str(resolved).startswith(str(cwd)):
        return resolved, f"Error: path '{path}' is outside the working directory."
    return resolved, None

def _run_bash(command: str, cwd: Path) -> str:
    try:
        r = subprocess.run(command, shell=True, cwd=str(cwd), capture_output=True, text=True,
                           timeout=BASH_TIMEOUT, env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1"})
        parts = [f"Exit code: {r.returncode}"]
        if r.stdout and r.stdout.strip(): parts.append(f"STDOUT:\n{r.stdout[:MAX_BASH_OUTPUT]}")
        if r.stderr and r.stderr.strip(): parts.append(f"STDERR:\n{r.stderr[:MAX_BASH_OUTPUT]}")
        return "\n".join(parts)
    except subprocess.TimeoutExpired: return f"Command timed out after {BASH_TIMEOUT}s."
    except Exception as e: return f"Error: {e}"

def _read_pdf(resolved: Path, display_path: str) -> str:
    try:
        import pymupdf
        doc = pymupdf.open(str(resolved)); parts = []
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if text: parts.append(f"--- Page {i+1} ---\n{text}")
        doc.close()
        result = "\n\n".join(parts) if parts else "(no text)"
        if len(result) > MAX_FILE_READ: result = result[:MAX_FILE_READ] + "\n...(truncated)"
        return f"File: {display_path} (PDF, {len(parts)} pages)\n{result}"
    except Exception as e: return f"Error reading PDF: {e}"

def _read_office(resolved: Path, display_path: str) -> str:
    suffix = resolved.suffix.lower()
    try:
        if suffix == ".docx":
            from docx import Document as D
            doc = D(str(resolved)); parts = [p.text for p in doc.paragraphs]
            for t in doc.tables:
                for row in t.rows: parts.append(" | ".join(c.text for c in row.cells))
            result = "\n".join(parts)
        elif suffix in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(str(resolved), data_only=True); parts = []
            for sn in wb.sheetnames:
                ws = wb[sn]; parts.append(f"[Sheet: {sn}]")
                for row in ws.iter_rows(values_only=True):
                    parts.append(" | ".join(str(c) if c is not None else "" for c in row))
            result = "\n".join(parts)
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(resolved)); parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"[Slide {i+1}]")
                for shape in slide.shapes:
                    if shape.has_text_frame: parts.append(shape.text_frame.text)
            result = "\n".join(parts)
        else: return f"Unsupported: {suffix}"
        if len(result) > MAX_FILE_READ: result = result[:MAX_FILE_READ] + "\n...(truncated)"
        return f"File: {display_path} ({suffix})\n{result}"
    except Exception as e: return f"Error reading '{display_path}': {e}"

def _read_file(path: str, cwd: Path, offset: int = 0, limit: int = 0) -> str:
    resolved, err = _check_path(path, cwd)
    if err: return err
    if not resolved.exists():
        avail = "\n".join(f"  {i.name}{'/' if i.is_dir() else ''}" for i in sorted(cwd.iterdir()) if not i.name.startswith(".")) or "(empty)"
        return f"Error: '{path}' not found.\nAvailable:\n{avail}"
    if resolved.is_dir(): return f"Error: '{path}' is a directory."
    if resolved.suffix.lower() == ".pdf": return _read_pdf(resolved, path)
    if resolved.suffix.lower() in (".docx", ".xlsx", ".xlsm", ".pptx"): return _read_office(resolved, path)
    try: content = resolved.read_text(encoding="utf-8", errors="strict")
    except (UnicodeDecodeError, ValueError): return f"Binary file ({resolved.stat().st_size} bytes)."
    lines = content.splitlines(keepends=True); total = len(lines)
    start = max(0, offset - 1) if offset > 0 else 0
    end = min(start + limit, total) if limit > 0 else total
    numbered = [f"{i:>6}\t{l.rstrip()}" for i, l in enumerate(lines[start:end], start=start+1)]
    result = "\n".join(numbered)
    if len(result) > MAX_FILE_READ: result = result[:MAX_FILE_READ] + f"\n...({total} lines)"
    header = f"File: {path} ({total} lines)"
    if start > 0 or end < total: header += f" [lines {start+1}-{end}]"
    return f"{header}\n{result}"

def _write_file(path: str, content: str, cwd: Path) -> str:
    resolved, err = _check_path(path, cwd)
    if err: return err
    try:
        resolved.parent.mkdir(parents=True, exist_ok=True); resolved.write_text(content)
        return f"Wrote {resolved.stat().st_size} bytes to {path}"
    except Exception as e: return f"Error: {e}"

def _edit_file(path: str, old_string: str, new_string: str, cwd: Path) -> str:
    resolved, err = _check_path(path, cwd)
    if err: return err
    if not resolved.exists(): return f"Error: '{path}' not found."
    try: content = resolved.read_text(encoding="utf-8")
    except: return "Error: binary file."
    count = content.count(old_string)
    if count == 0: return f"Error: not found.\nPreview:\n{content[:2000]}"
    if count > 1: return f"Error: found {count} times. Add context."
    resolved.write_text(content.replace(old_string, new_string, 1))
    return f"Edited {path}."

def _list_dir(path: str, cwd: Path) -> str:
    resolved, err = _check_path(path, cwd)
    if err: return err
    if not resolved.exists(): return f"Error: '{path}' doesn't exist."
    if not resolved.is_dir(): return f"Error: not a directory."
    entries = []
    for item in sorted(resolved.iterdir()):
        if item.name.startswith("."): continue
        if item.is_file():
            s = item.stat().st_size
            sz = f"{s}B" if s < 1024 else f"{s/1024:.1f}KB" if s < 1048576 else f"{s/1048576:.1f}MB"
            entries.append(f"  {item.name}  ({sz})")
        elif item.is_dir(): entries.append(f"  {item.name}/")
    return "\n".join(entries) if entries else "(empty)"

def _grep(pattern: str, path: str, cwd: Path, include: str = "") -> str:
    resolved, err = _check_path(path, cwd)
    if err: return err
    if not resolved.exists(): return f"Error: '{path}' doesn't exist."
    try: regex = re.compile(pattern)
    except re.error as e: return f"Error: invalid regex: {e}"
    matches: list[str] = []
    def _s(fp):
        try: text = fp.read_text(encoding="utf-8", errors="replace")
        except: return
        rel = str(fp.relative_to(cwd))
        for i, line in enumerate(text.splitlines(), 1):
            if regex.search(line):
                matches.append(f"{rel}:{i}: {line.rstrip()[:200]}")
                if len(matches) >= MAX_GREP_RESULTS: return
    if resolved.is_file(): _s(resolved)
    else:
        for fp in sorted(resolved.rglob("*")):
            if len(matches) >= MAX_GREP_RESULTS: break
            if not fp.is_file() or fp.name.startswith("."): continue
            if include and not fnmatch.fnmatch(fp.name, include): continue
            _s(fp)
    if not matches: return f"No matches for '{pattern}'."
    h = f"Found {len(matches)} match(es)"
    if len(matches) >= MAX_GREP_RESULTS: h += f" (limit {MAX_GREP_RESULTS})"
    return f"{h}:\n" + "\n".join(matches)

def _glob_files(pattern: str, cwd: Path) -> str:
    try: matches = sorted(cwd.glob(pattern), key=lambda p: p.stat().st_mtime, reverse=True)
    except: matches = []
    matches = [m for m in matches if not any(p.startswith(".") for p in m.relative_to(cwd).parts)]
    if not matches: return f"No files matching '{pattern}'."
    results = [f"  {str(m.relative_to(cwd))}  ({m.stat().st_size}B)" if m.is_file()
               else f"  {str(m.relative_to(cwd))}/" for m in matches[:MAX_GLOB_RESULTS]]
    h = f"Found {len(matches)} match(es)"
    if len(matches) > MAX_GLOB_RESULTS: h += f" (showing {MAX_GLOB_RESULTS})"
    return f"{h}:\n" + "\n".join(results)

def _done(summary: str) -> str:
    return f"TASK COMPLETE: {summary}"

def _execute_tool(name: str, args: dict[str, Any], cwd: Path) -> str:
    """Dispatch to the 8 standard tools."""
    try:
        if name == "bash": return _run_bash(args.get("command", ""), cwd)
        elif name == "read_file": return _read_file(args.get("path", ""), cwd, int(args.get("offset", 0)), int(args.get("limit", 0)))
        elif name == "write_file": return _write_file(args.get("path", ""), args.get("content", ""), cwd)
        elif name == "edit_file": return _edit_file(args.get("path", ""), args.get("old_string", ""), args.get("new_string", ""), cwd)
        elif name == "list_dir": return _list_dir(args.get("path", "."), cwd)
        elif name == "grep": return _grep(args.get("pattern", ""), args.get("path", "."), cwd, args.get("include", ""))
        elif name == "glob": return _glob_files(args.get("pattern", ""), cwd)
        elif name == "done": return _done(args.get("summary", ""))
        else: return f"Unknown tool: {name}"
    except Exception as e: return f"Tool error ({name}): {e}"

# ── Agent loop (Gemini-specific) ───────────────────────────────────────────

class CustomAgent(BaseAgent):
    """Self-contained agent using Gemini API with 8 standard tools."""

    def name(self) -> str: return "custom-gemini"

    async def run(self, prompt: str, cwd: Path) -> AgentResult:
        cwd = cwd.resolve(); cwd.mkdir(parents=True, exist_ok=True)
        tc_log: list[dict[str, Any]] = []; msg_log: list[dict[str, Any]] = []; resp: list[str] = []
        try:
            return await asyncio.to_thread(self._loop, prompt, cwd, tc_log, msg_log, resp)
        except Exception as e:
            logger.error("[gemini] Error: %s", e)
            return AgentResult(response=f"(Error: {e})", tool_calls=tc_log, messages=msg_log)

    def _loop(self, prompt, cwd, tc_log, msg_log, resp) -> AgentResult:
        api_key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
        if not api_key: return AgentResult(error="GEMINI_API_KEY or GOOGLE_API_KEY required.")
        client = genai.Client(api_key=api_key)
        model = self._model or DEFAULT_MODEL
        max_iters = self._max_turns if self._max_turns != 10 else MAX_ITERATIONS
        config = GenerateContentConfig(
            system_instruction=SYSTEM_PROMPT.format(cwd=cwd),
            temperature=TEMPERATURE, max_output_tokens=MAX_OUTPUT_TOKENS,
            tools=[TOOL_DECLARATIONS],
        )
        contents: list[Content | str] = [prompt]
        last_calls: list[str] = []; nudge = 0

        for it in range(max_iters):
            response = _retry(client, model, contents, config)
            if response is None: resp.append("API errors."); break
            fcs = _extract_fcs(response)

            if not fcs:
                text = _get_text(response); resp.append(text)
                msg_log.append({"role": "assistant", "type": "text", "content": text[:2000]})
                break

            for part in _get_parts(response):
                if part.text: resp.append(part.text); msg_log.append({"role": "assistant", "type": "text", "content": part.text[:2000]})

            tool_parts: list[Part] = []
            done_called = False
            for fc in fcs:
                try: args = dict(fc.args) if fc.args else {}
                except: args = {}
                result = _execute_tool(fc.name, args, cwd)
                tc_log.append({"tool": fc.name, "input": str(args)[:500]})
                msg_log.append({"role": "assistant", "type": "tool_use", "tool": fc.name, "input": str(args)[:1000]})
                msg_log.append({"role": "tool", "type": "tool_result", "content": result[:2000]})
                tool_parts.append(Part.from_function_response(name=fc.name, response={"result": result}))
                last_calls.append(f"{fc.name}:{hash(str(args))}")
                if fc.name == "done": done_called = True; resp.append(result)

            try: mc = response.candidates[0].content if response.candidates else None
            except: mc = None
            if mc: contents.append(mc)
            contents.append(Content(parts=tool_parts))

            if done_called: break

            if len(last_calls) >= 6 and last_calls[-3:] == last_calls[-6:-3]:
                nudge += 1
                if nudge >= 3: resp.append("(Terminated: stuck)"); break
                contents.append("You are repeating actions. Try different approach or call done.")

        return AgentResult(response="\n".join(resp), tool_calls=tc_log, messages=msg_log)

# ── Helpers (Gemini-specific) ──────────────────────────────────────────────

def _retry(client, model, contents, config, retries=3):
    for i in range(retries + 1):
        try: return client.models.generate_content(model=model, contents=contents, config=config)
        except Exception as e:
            if i == retries: logger.error("[gemini] API failed: %s", e); return None
            time.sleep(2 ** (i + 1))
    return None

def _extract_fcs(response):
    if not response.candidates: return []
    c = response.candidates[0].content
    if not c or not c.parts: return []
    return [p.function_call for p in c.parts if p.function_call]

def _get_text(response):
    try: return response.text or ""
    except: pass
    if response.candidates:
        c = response.candidates[0].content
        if c and c.parts: return "\n".join(p.text for p in c.parts if p.text)
    return ""

def _get_parts(response):
    if not response.candidates: return []
    c = response.candidates[0].content
    return c.parts if c and c.parts else []
