"""Evaluation runner — sends GDPval tasks to agent backends and scores responses."""

from __future__ import annotations

import asyncio
import json
import logging
import shutil
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from src.data.types import Sample
from src.eval.agents.base import BaseAgent
from src.eval.evaluators.base import BaseEvaluator, EvalResult
from src.eval.evaluators.gdpval import GDPvalEvaluator
from src.eval.evaluators.gdpval_judge import GDPvalJudgeEvaluator

logger = logging.getLogger(__name__)


@dataclass
class TaskTrace:
    """Full trace of a single task execution."""

    task_id: str
    benchmark: str
    prompt: str
    agent_name: str = ""
    response: str = ""
    tool_calls: list[dict[str, Any]] = field(default_factory=list)
    messages: list[dict[str, Any]] = field(default_factory=list)
    eval_result: EvalResult | None = None
    duration_s: float = 0.0
    error: str | None = None
    workspace_dir: str = ""


@dataclass
class BatchResult:
    """Result of running a batch of tasks."""

    traces: list[TaskTrace] = field(default_factory=list)
    total_duration_s: float = 0.0

    @property
    def avg_score(self) -> float:
        if not self.traces:
            return 0.0
        total = 0.0
        scorable = 0
        for t in self.traces:
            if t.eval_result and t.eval_result.max_score > 0:
                total += t.eval_result.normalized_score
                scorable += 1
        return total / scorable if scorable > 0 else 0.0

    @property
    def num_completed(self) -> int:
        return sum(1 for t in self.traces if t.error is None)

    @property
    def num_errors(self) -> int:
        return sum(1 for t in self.traces if t.error is not None)


class GDPvalRunner:
    """Runs GDPval tasks through any agent backend and evaluates responses.

    Supports Claude Code, Gemini CLI, and Codex CLI via pluggable agents.

    Usage:
        from src.eval.agents import ClaudeCodeAgent
        agent = ClaudeCodeAgent(max_turns=10)
        runner = GDPvalRunner(agent=agent)
        result = await runner.run_batch(samples[:5])
    """

    def __init__(
        self,
        agent: BaseAgent,
        working_dir: str | Path | None = None,
        use_judge: bool = True,
        judge_model: str | None = None,
    ) -> None:
        self._agent = agent
        self._base_working_dir = Path(working_dir) if working_dir else Path("results/workspace")

        if use_judge:
            kwargs: dict[str, Any] = {}
            if judge_model:
                kwargs["model"] = judge_model
            self._evaluator: BaseEvaluator = GDPvalJudgeEvaluator(**kwargs)
            logger.info("Using Gemini judge evaluator (%s)", self._evaluator.name())
        else:
            self._evaluator = GDPvalEvaluator()
            logger.info("Using keyword-based evaluator (heuristic)")

    # File extensions to skip (generated scripts, temp files, etc.)
    _SKIP_PATTERNS = {".pyc", ".o", ".so", ".dylib", ".DS_Store"}

    @staticmethod
    def _extract_workspace_files(task_dir: Path) -> dict[str, str]:
        """Extract text content from all deliverable files in the workspace.

        Supports: .docx, .xlsx, .xls, .pdf, .pptx, .csv, .txt, .md,
        .json, .jsonl, .html, .xml, .rtf
        """
        contents: dict[str, str] = {}
        for path in task_dir.iterdir():
            if path.name.startswith("."):
                continue
            if path.suffix in GDPvalRunner._SKIP_PATTERNS:
                continue
            if path.suffix == ".py":
                continue
            try:
                text = GDPvalRunner._extract_file(path)
                if text:
                    contents[path.name] = text
            except Exception as e:
                logger.warning("Could not extract %s: %s", path.name, e)
        return contents

    @staticmethod
    def _extract_file(path: Path) -> str | None:
        """Extract text from a single file based on its extension."""
        suffix = path.suffix.lower()

        if suffix == ".docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(str(path))
            parts: list[str] = []
            for para in doc.paragraphs:
                parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
            return "\n".join(parts)

        if suffix in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(str(path), data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    row_str = " | ".join(str(c) if c is not None else "" for c in row)
                    parts.append(row_str)
            return "\n".join(parts)

        if suffix == ".xls":
            import xlrd
            wb = xlrd.open_workbook(str(path))
            parts = []
            for sheet in wb.sheets():
                parts.append(f"[Sheet: {sheet.name}]")
                for row_idx in range(sheet.nrows):
                    row_vals = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                    parts.append(" | ".join(row_vals))
            return "\n".join(parts)

        if suffix == ".pdf":
            import pymupdf
            doc = pymupdf.open(str(path))
            parts = []
            for page in doc:
                parts.append(page.get_text())
            doc.close()
            return "\n".join(parts)

        if suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"[Slide {i + 1}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            parts.append(" | ".join(cell.text for cell in row.cells))
            return "\n".join(parts)

        if suffix in (".csv", ".txt", ".md", ".json", ".jsonl", ".html", ".xml",
                       ".rtf", ".log", ".yaml", ".yml", ".toml", ".ini", ".cfg"):
            return path.read_text(errors="replace")[:50000]

        return None

    @staticmethod
    def _copy_reference_files(sample: Sample, task_dir: Path) -> int:
        """Copy reference files into the task workspace. Returns count copied."""
        ref_files = sample.metadata.get("reference_files", [])
        if not hasattr(ref_files, '__len__') or len(ref_files) == 0:
            return 0
        gdpval_dir = Path("data/raw/gdpval")
        copied = 0
        for rel_path in ref_files:
            src = gdpval_dir / rel_path
            if src.exists():
                dst = task_dir / src.name
                shutil.copy2(str(src), str(dst))
                copied += 1
                logger.debug("Copied reference file: %s → %s", src.name, dst)
            else:
                logger.warning("Reference file not found: %s", src)
        return copied

    async def run_single(self, sample: Sample) -> TaskTrace:
        """Run a single GDPval task and evaluate the response."""
        task_dir = self._base_working_dir / sample.id[:12]
        task_dir.mkdir(parents=True, exist_ok=True)

        # Copy reference files into workspace
        n_copied = self._copy_reference_files(sample, task_dir)
        if n_copied:
            logger.debug("[files] Copied %d reference files to workspace", n_copied)

        occupation = sample.metadata.get("occupation", "unknown")
        logger.debug(
            "[%s] Starting task %s (%s) — workspace: %s",
            self._agent.name(), sample.id[:12], occupation, task_dir,
        )

        trace = TaskTrace(
            task_id=sample.id,
            benchmark=sample.benchmark,
            prompt=sample.prompt,
            agent_name=self._agent.name(),
            workspace_dir=str(task_dir),
        )
        start = time.monotonic()

        try:
            agent_result = await self._agent.run(prompt=sample.prompt, cwd=task_dir)

            if agent_result.error:
                trace.error = agent_result.error
            else:
                trace.response = agent_result.response
                trace.tool_calls = agent_result.tool_calls
                trace.messages = agent_result.messages

                # Extract content from deliverable files in workspace
                file_contents = self._extract_workspace_files(task_dir)
                if file_contents:
                    logger.debug(
                        "[%s] Extracted content from %d workspace files",
                        self._agent.name(), len(file_contents),
                    )
                    trace.response += "\n\n--- DELIVERABLE FILE CONTENTS ---\n"
                    for fname, content in file_contents.items():
                        trace.response += f"\n=== {fname} ===\n{content}\n"

                logger.debug(
                    "[%s] Response: %d chars, %d tool calls",
                    self._agent.name(), len(trace.response), len(trace.tool_calls),
                )

                # Evaluate against rubric
                logger.debug("[eval] Scoring task %s with %s...", sample.id[:12], self._evaluator.name())
                eval_start = time.monotonic()

                trace.eval_result = await self._evaluator.evaluate(
                    task_id=sample.id,
                    prompt=sample.prompt,
                    response=trace.response,
                    reference=sample.reference,
                    rubric_json=sample.metadata.get("rubric_json", ""),
                )

                eval_duration = time.monotonic() - eval_start
                if trace.eval_result:
                    logger.info(
                        "[eval] Task %s scored: %.1f/%.1f (%.0f%%) — %d/%d criteria (%.1fs)",
                        sample.id[:12],
                        trace.eval_result.score,
                        trace.eval_result.max_score,
                        trace.eval_result.normalized_score * 100,
                        trace.eval_result.metadata.get("criteria_met", 0),
                        trace.eval_result.metadata.get("num_criteria", 0),
                        eval_duration,
                    )

        except Exception as e:
            trace.error = f"{type(e).__name__}: {e}"
            logger.error("[error] Task %s failed: %s", sample.id[:12], trace.error)

        trace.duration_s = time.monotonic() - start
        logger.info("[done] Task %s total time: %.1fs", sample.id[:12], trace.duration_s)
        return trace

    async def run_batch(
        self,
        samples: list[Sample],
        concurrency: int = 1,
        progress_callback: Any = None,
    ) -> BatchResult:
        """Run a batch of GDPval tasks."""
        start = time.monotonic()
        result = BatchResult()

        logger.info(
            "[batch] Starting: %d tasks, concurrency=%d, agent=%s, evaluator=%s",
            len(samples), concurrency, self._agent.name(), self._evaluator.name(),
        )

        if concurrency <= 1:
            for i, sample in enumerate(samples):
                logger.info(
                    "━━━ Task %d/%d ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
                    i + 1, len(samples),
                )
                trace = await self.run_single(sample)
                result.traces.append(trace)

                scored = [t.eval_result for t in result.traces if t.eval_result and t.eval_result.max_score > 0]
                running_avg = sum(r.normalized_score for r in scored) / len(scored) if scored else 0
                logger.info(
                    "[batch] Progress: %d/%d done — running avg: %.1f%%",
                    i + 1, len(samples), running_avg * 100,
                )
                if progress_callback:
                    progress_callback(i + 1, len(samples), trace)
        else:
            sem = asyncio.Semaphore(concurrency)
            completed = 0
            scores: list[float] = []
            errors = 0

            async def run_with_sem(sample: Sample) -> TaskTrace:
                nonlocal completed, errors
                async with sem:
                    trace = await self.run_single(sample)
                    completed += 1
                    if trace.error:
                        errors += 1
                    elif trace.eval_result and trace.eval_result.max_score > 0:
                        scores.append(trace.eval_result.normalized_score)
                    if progress_callback:
                        progress_callback(completed, len(samples), trace)
                    return trace

            traces = await asyncio.gather(
                *[run_with_sem(s) for s in samples],
                return_exceptions=False,
            )
            result.traces = list(traces)

            elapsed = time.monotonic() - start
            avg = sum(scores) / len(scores) * 100 if scores else 0
            logger.info(
                "[batch] Done: %d tasks in %.0fs — avg %.1f%%, %d errors",
                len(samples), elapsed, avg, errors,
            )

        result.total_duration_s = time.monotonic() - start
        return result

    @staticmethod
    def save_results(batch: BatchResult, output_dir: Path) -> tuple[Path, Path]:
        """Save batch results as traces.json + eval.json."""
        output_dir.mkdir(parents=True, exist_ok=True)

        # --- traces.json ---
        trace_data = {
            "num_tasks": len(batch.traces),
            "total_duration_s": batch.total_duration_s,
            "traces": [],
        }
        for trace in batch.traces:
            trace_data["traces"].append({
                "task_id": trace.task_id,
                "benchmark": trace.benchmark,
                "agent": trace.agent_name,
                "prompt": trace.prompt,
                "response": trace.response,
                "tool_calls": trace.tool_calls,
                "messages": trace.messages,
                "duration_s": trace.duration_s,
                "error": trace.error,
                "workspace_dir": trace.workspace_dir,
            })

        trace_path = output_dir / "traces.json"
        trace_path.write_text(json.dumps(trace_data, indent=2, default=str))
        logger.info("Traces saved to %s", trace_path)

        # --- eval.json ---
        eval_data: dict[str, Any] = {
            "summary": {
                "agent": batch.traces[0].agent_name if batch.traces else "",
                "avg_score": batch.avg_score,
                "num_tasks": len(batch.traces),
                "num_completed": batch.num_completed,
                "num_errors": batch.num_errors,
                "total_duration_s": batch.total_duration_s,
            },
            "scores": [],
        }
        for trace in batch.traces:
            entry: dict[str, Any] = {
                "task_id": trace.task_id,
                "error": trace.error,
            }
            if trace.eval_result:
                entry["score"] = trace.eval_result.score
                entry["max_score"] = trace.eval_result.max_score
                entry["normalized_score"] = trace.eval_result.normalized_score
                entry["rubric_breakdown"] = trace.eval_result.rubric_breakdown
                entry["criteria_met"] = trace.eval_result.metadata.get("criteria_met", 0)
                entry["num_criteria"] = trace.eval_result.metadata.get("num_criteria", 0)
            eval_data["scores"].append(entry)

        eval_path = output_dir / "eval.json"
        eval_path.write_text(json.dumps(eval_data, indent=2, default=str))
        logger.info("Eval results saved to %s", eval_path)

        return trace_path, eval_path
