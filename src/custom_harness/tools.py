"""Tool declarations and implementations for the custom harness agent.

Provides 7 tools modeled after Claude Code and Codex CLI:
    bash      — Execute shell commands
    read_file — Read file contents (with offset/limit)
    write_file— Create or overwrite files
    edit_file — Find-and-replace edits
    list_dir  — List directory contents
    grep      — Search file contents by regex
    glob      — Find files by name pattern
"""

from __future__ import annotations

import fnmatch
import os
import re
import subprocess
from pathlib import Path
from typing import Any

from google.genai.types import FunctionDeclaration, Tool

from src.custom_harness.config import (
    BASH_TIMEOUT,
    MAX_BASH_OUTPUT,
    MAX_FILE_READ,
    MAX_GLOB_RESULTS,
    MAX_GREP_RESULTS,
)

# ── Tool declarations (Gemini FunctionDeclaration format) ───────────────────

TOOL_DECLARATIONS = Tool(function_declarations=[
    FunctionDeclaration(
        name="bash",
        description=(
            "Execute a shell command in the working directory. "
            "Returns stdout, stderr, and exit code. Use for running scripts, "
            "installing packages, git operations, and any system commands."
        ),
        parameters={
            "type": "object",
            "properties": {
                "command": {
                    "type": "string",
                    "description": "The shell command to execute.",
                },
            },
            "required": ["command"],
        },
    ),
    FunctionDeclaration(
        name="read_file",
        description=(
            "Read the contents of a file. Supports text files (with line numbers), "
            "PDFs (extracts text from all pages), and Office documents (.docx, .xlsx, .pptx). "
            "For large files, use offset and limit to read specific line ranges."
        ),
        parameters={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the file (relative to working directory).",
                },
                "offset": {
                    "type": "integer",
                    "description": "Line number to start reading from (1-based). Optional.",
                },
                "limit": {
                    "type": "integer",
                    "description": "Maximum number of lines to read. Optional.",
                },
            },
            "required": ["path"],
        },
    ),
    FunctionDeclaration(
        name="write_file",
        description=(
            "Create a new file or completely overwrite an existing file. "
            "Use for creating scripts, config files, text documents, etc. "
            "For targeted edits to existing files, use edit_file instead."
        ),
        parameters={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path for the file (relative to working directory).",
                },
                "content": {
                    "type": "string",
                    "description": "The complete file content to write.",
                },
            },
            "required": ["path", "content"],
        },
    ),
    FunctionDeclaration(
        name="edit_file",
        description=(
            "Make targeted find-and-replace edits in an existing file. "
            "The old_string must match exactly (including whitespace/indentation). "
            "Use this for surgical edits; use write_file for full rewrites."
        ),
        parameters={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the file (relative to working directory).",
                },
                "old_string": {
                    "type": "string",
                    "description": "The exact text to find in the file.",
                },
                "new_string": {
                    "type": "string",
                    "description": "The text to replace it with.",
                },
            },
            "required": ["path", "old_string", "new_string"],
        },
    ),
    FunctionDeclaration(
        name="list_dir",
        description=(
            "List files and subdirectories in a directory. "
            "Returns names with file sizes. Useful for understanding "
            "workspace contents and finding files."
        ),
        parameters={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Directory path (relative to working directory). Defaults to '.'.",
                },
            },
        },
    ),
    FunctionDeclaration(
        name="grep",
        description=(
            "Search file contents using a regex pattern. "
            "Returns matching lines with file paths and line numbers. "
            "Searches recursively through the specified path."
        ),
        parameters={
            "type": "object",
            "properties": {
                "pattern": {
                    "type": "string",
                    "description": "Regular expression pattern to search for.",
                },
                "path": {
                    "type": "string",
                    "description": "File or directory to search in (relative). Defaults to '.'.",
                },
                "include": {
                    "type": "string",
                    "description": "Glob pattern to filter files (e.g. '*.py', '*.txt'). Optional.",
                },
            },
            "required": ["pattern"],
        },
    ),
    FunctionDeclaration(
        name="glob",
        description=(
            "Find files matching a glob pattern. Returns file paths sorted "
            "by modification time. Supports ** for recursive matching."
        ),
        parameters={
            "type": "object",
            "properties": {
                "pattern": {
                    "type": "string",
                    "description": "Glob pattern (e.g. '**/*.py', '*.xlsx', 'src/**/*.ts').",
                },
            },
            "required": ["pattern"],
        },
    ),
])


# ── Path safety ─────────────────────────────────────────────────────────────

def _check_path(path: str, cwd: Path) -> tuple[Path, str | None]:
    """Resolve a path relative to cwd and verify it stays within.

    Returns (resolved_path, error_message_or_None).
    """
    resolved = (cwd / path).resolve()
    if not str(resolved).startswith(str(cwd)):
        return resolved, f"Error: path '{path}' is outside the working directory."
    return resolved, None


def _list_available_files(cwd: Path) -> str:
    """List files in workspace root — used in error messages."""
    entries = []
    for item in sorted(cwd.iterdir()):
        if item.name.startswith("."):
            continue
        if item.is_file():
            entries.append(f"  {item.name}")
        elif item.is_dir():
            entries.append(f"  {item.name}/")
    return "\n".join(entries) if entries else "(empty directory)"


# ── Tool implementations ───────────────────────────────────────────────────

def run_bash(command: str, cwd: Path) -> str:
    """Execute a shell command in the workspace."""
    try:
        result = subprocess.run(
            command,
            shell=True,
            cwd=str(cwd),
            capture_output=True,
            text=True,
            timeout=BASH_TIMEOUT,
            env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1"},
        )
        stdout = result.stdout[:MAX_BASH_OUTPUT] if result.stdout else ""
        stderr = result.stderr[:MAX_BASH_OUTPUT] if result.stderr else ""
        parts = [f"Exit code: {result.returncode}"]
        if stdout.strip():
            parts.append(f"STDOUT:\n{stdout}")
        if stderr.strip():
            parts.append(f"STDERR:\n{stderr}")
        return "\n".join(parts)
    except subprocess.TimeoutExpired:
        return f"Command timed out after {BASH_TIMEOUT}s. Consider breaking into smaller steps."
    except Exception as e:
        return f"Error executing command: {e}"


def _read_pdf(resolved: Path, display_path: str) -> str:
    """Extract text from a PDF file using pymupdf."""
    try:
        import pymupdf
        doc = pymupdf.open(str(resolved))
        parts = []
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if text:
                parts.append(f"--- Page {i + 1} ---\n{text}")
        doc.close()
        result = "\n\n".join(parts) if parts else "(PDF contains no extractable text)"
        if len(result) > MAX_FILE_READ:
            result = result[:MAX_FILE_READ] + "\n... (truncated)"
        return f"File: {display_path} (PDF, {len(parts)} pages)\n{result}"
    except Exception as e:
        return f"Error reading PDF '{display_path}': {e}"


def _read_office_file(resolved: Path, display_path: str) -> str:
    """Extract text from Office documents (.docx, .xlsx, .pptx)."""
    suffix = resolved.suffix.lower()
    try:
        if suffix == ".docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(str(resolved))
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
            result = "\n".join(parts)

        elif suffix in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(str(resolved), data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    row_str = " | ".join(str(c) if c is not None else "" for c in row)
                    parts.append(row_str)
            result = "\n".join(parts)

        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(resolved))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"[Slide {i + 1}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            result = "\n".join(parts)
        else:
            return f"Unsupported file type: {suffix}"

        if len(result) > MAX_FILE_READ:
            result = result[:MAX_FILE_READ] + "\n... (truncated)"
        return f"File: {display_path} ({suffix})\n{result}"

    except Exception as e:
        return f"Error reading '{display_path}': {e}"


def read_file(path: str, cwd: Path, offset: int = 0, limit: int = 0) -> str:
    """Read a file with optional offset/limit (line-based)."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    if not resolved.exists():
        available = _list_available_files(cwd)
        return f"Error: file '{path}' not found.\nAvailable files:\n{available}"
    if resolved.is_dir():
        return f"Error: '{path}' is a directory. Use list_dir instead."

    # Handle PDFs natively using pymupdf
    if resolved.suffix.lower() == ".pdf":
        return _read_pdf(resolved, path)

    # Handle Office documents via Python libs
    if resolved.suffix.lower() in (".docx", ".xlsx", ".xlsm", ".pptx"):
        return _read_office_file(resolved, path)

    try:
        content = resolved.read_text(encoding="utf-8", errors="strict")
    except (UnicodeDecodeError, ValueError):
        size = resolved.stat().st_size
        return f"Binary file ({size} bytes). Use bash with Python to inspect binary files."

    lines = content.splitlines(keepends=True)
    total_lines = len(lines)

    start = max(0, offset - 1) if offset > 0 else 0
    end = min(start + limit, total_lines) if limit > 0 else total_lines
    selected = lines[start:end]

    numbered = []
    for i, line in enumerate(selected, start=start + 1):
        numbered.append(f"{i:>6}\t{line.rstrip()}")
    result = "\n".join(numbered)

    if len(result) > MAX_FILE_READ:
        result = result[:MAX_FILE_READ] + f"\n... (truncated, {total_lines} total lines)"

    header = f"File: {path} ({total_lines} lines)"
    if start > 0 or end < total_lines:
        header += f" [showing lines {start + 1}-{end}]"
    return f"{header}\n{result}"


def write_file(path: str, content: str, cwd: Path) -> str:
    """Create or overwrite a file in the workspace."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    try:
        resolved.parent.mkdir(parents=True, exist_ok=True)
        resolved.write_text(content)
        size = resolved.stat().st_size
        return f"Successfully wrote {size} bytes to {path}"
    except Exception as e:
        return f"Error writing file: {e}"


def edit_file(path: str, old_string: str, new_string: str, cwd: Path) -> str:
    """Find-and-replace edit in a file."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    if not resolved.exists():
        return f"Error: file '{path}' not found."
    try:
        content = resolved.read_text(encoding="utf-8")
    except (UnicodeDecodeError, ValueError):
        return f"Error: '{path}' is a binary file and cannot be edited with edit_file."

    count = content.count(old_string)
    if count == 0:
        preview = content[:2000] if len(content) > 2000 else content
        return (
            f"Error: old_string not found in '{path}'. "
            f"Make sure the text matches exactly (including whitespace).\n"
            f"File preview:\n{preview}"
        )
    if count > 1:
        return (
            f"Error: old_string found {count} times in '{path}'. "
            f"Provide more surrounding context to make the match unique."
        )

    new_content = content.replace(old_string, new_string, 1)
    resolved.write_text(new_content)

    new_lines = new_content.splitlines()
    edit_start = content.index(old_string)
    line_num = content[:edit_start].count("\n")
    ctx_start = max(0, line_num - 2)
    ctx_end = min(len(new_lines), line_num + new_string.count("\n") + 3)
    context = "\n".join(
        f"{i + 1:>6}\t{new_lines[i]}" for i in range(ctx_start, ctx_end)
    )
    return f"Successfully edited {path}.\nUpdated region:\n{context}"


def list_dir(path: str, cwd: Path) -> str:
    """List contents of a directory."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    if not resolved.exists():
        return f"Error: directory '{path}' does not exist."
    if not resolved.is_dir():
        return f"Error: '{path}' is not a directory."

    entries = []
    for item in sorted(resolved.iterdir()):
        if item.name.startswith("."):
            continue
        if item.is_file():
            size = item.stat().st_size
            if size < 1024:
                size_str = f"{size}B"
            elif size < 1024 * 1024:
                size_str = f"{size / 1024:.1f}KB"
            else:
                size_str = f"{size / (1024 * 1024):.1f}MB"
            entries.append(f"  {item.name}  ({size_str})")
        elif item.is_dir():
            entries.append(f"  {item.name}/")
    return "\n".join(entries) if entries else "(empty directory)"


def grep_files(pattern: str, path: str, cwd: Path, include: str = "") -> str:
    """Search file contents using regex."""
    resolved, err = _check_path(path, cwd)
    if err:
        return err
    if not resolved.exists():
        return f"Error: path '{path}' does not exist."

    try:
        regex = re.compile(pattern)
    except re.error as e:
        return f"Error: invalid regex pattern: {e}"

    matches: list[str] = []

    def _search_file(fpath: Path) -> None:
        try:
            text = fpath.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return
        rel = str(fpath.relative_to(cwd))
        for i, line in enumerate(text.splitlines(), 1):
            if regex.search(line):
                matches.append(f"{rel}:{i}: {line.rstrip()[:200]}")
                if len(matches) >= MAX_GREP_RESULTS:
                    return

    if resolved.is_file():
        _search_file(resolved)
    else:
        for fpath in sorted(resolved.rglob("*")):
            if len(matches) >= MAX_GREP_RESULTS:
                break
            if not fpath.is_file() or fpath.name.startswith("."):
                continue
            if include and not fnmatch.fnmatch(fpath.name, include):
                continue
            _search_file(fpath)

    if not matches:
        return f"No matches found for pattern '{pattern}' in {path}."
    header = f"Found {len(matches)} match(es)"
    if len(matches) >= MAX_GREP_RESULTS:
        header += f" (limited to {MAX_GREP_RESULTS})"
    return f"{header}:\n" + "\n".join(matches)


def glob_files(pattern: str, cwd: Path) -> str:
    """Find files matching a glob pattern."""
    try:
        matches = sorted(
            cwd.glob(pattern),
            key=lambda p: p.stat().st_mtime if p.exists() else 0,
            reverse=True,
        )
    except Exception as e:
        return f"Error: invalid glob pattern: {e}"

    # Filter hidden files
    matches = [
        m for m in matches
        if not any(part.startswith(".") for part in m.relative_to(cwd).parts)
    ]
    if not matches:
        return f"No files matching '{pattern}'."

    results = []
    for m in matches[:MAX_GLOB_RESULTS]:
        rel = str(m.relative_to(cwd))
        if m.is_file():
            results.append(f"  {rel}  ({m.stat().st_size} bytes)")
        else:
            results.append(f"  {rel}/")
    header = f"Found {len(matches)} match(es)"
    if len(matches) > MAX_GLOB_RESULTS:
        header += f" (showing first {MAX_GLOB_RESULTS})"
    return f"{header}:\n" + "\n".join(results)


# ── Dispatcher ──────────────────────────────────────────────────────────────

def execute_tool(name: str, args: dict[str, Any], cwd: Path) -> str:
    """Dispatch a tool call to its implementation."""
    try:
        if name == "bash":
            return run_bash(args.get("command", ""), cwd)
        elif name == "read_file":
            return read_file(
                args.get("path", ""), cwd,
                offset=int(args.get("offset", 0)),
                limit=int(args.get("limit", 0)),
            )
        elif name == "write_file":
            return write_file(args.get("path", ""), args.get("content", ""), cwd)
        elif name == "edit_file":
            return edit_file(
                args.get("path", ""), args.get("old_string", ""),
                args.get("new_string", ""), cwd,
            )
        elif name == "list_dir":
            return list_dir(args.get("path", "."), cwd)
        elif name == "grep":
            return grep_files(
                args.get("pattern", ""), args.get("path", "."),
                cwd, include=args.get("include", ""),
            )
        elif name == "glob":
            return glob_files(args.get("pattern", ""), cwd)
        else:
            return f"Unknown tool: {name}"
    except Exception as e:
        return f"Tool error ({name}): {e}"
